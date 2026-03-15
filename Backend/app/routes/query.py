from fastapi import APIRouter, UploadFile, File
from fastapi.responses import Response
from pydantic import BaseModel
import urllib.request
from html.parser import HTMLParser
import io
import PyPDF2
import docx
import pptx

from app.services.embedding_service import collection
from app.services.gemini_service import answer_question, bullet_points, create_flow, generate_tags, extract_main_content, translate_content

class QueryRequest(BaseModel):
    question: str

class QueryResponse(BaseModel):
    answer: str

class ActionRequest(BaseModel):
    content: str

class UrlRequest(BaseModel):
    url: str

class TranslateRequest(BaseModel):
    content: str
    target_language: str

router = APIRouter()

@router.post("/ask", response_model=QueryResponse)
def ask_question(request: QueryRequest):
    results = collection.query(query_texts=[request.question], n_results=3)
    docs_list = results.get("documents")
    docs = docs_list[0] if docs_list and len(docs_list) > 0 else []
    context = "\n".join(docs) if docs else ""
    answer = answer_question(context, request.question)
    return {"answer": answer}

@router.post("/bullets")
def get_bullets(request: ActionRequest):
    return {"result": bullet_points(request.content)}

@router.post("/flow")
def get_flow(request: ActionRequest):
    return {"result": create_flow(request.content)}

@router.post("/tags")
def get_tags(request: ActionRequest):
    return {"result": generate_tags(request.content)}

@router.post("/translate")
def translate_text(request: TranslateRequest):
    return {"result": translate_content(request.content, request.target_language)}

class TextParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.text = []
        self.hide_output = False
        self.ignore_tags = {'script', 'style', 'head'}

    def handle_starttag(self, tag, attrs):
        if tag in self.ignore_tags:
            self.hide_output = True

    def handle_endtag(self, tag):
        if tag in self.ignore_tags:
            self.hide_output = False

    def handle_data(self, data):
        if not self.hide_output and data.strip():
            self.text.append(data.strip())
            
    def get_text(self):
        return ' '.join(self.text)

import requests

@router.post("/scrape")
def scrape_url(request: UrlRequest):
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7',
            'Accept-Language': 'en-US,en;q=0.9',
        }
        
        response = requests.get(request.url, headers=headers, timeout=10)
        response.raise_for_status()
        html = response.text
        
        parser = TextParser()
        parser.feed(html)
        # Extracted text limited to ensure it doesn't break limits instantly
        raw_text = parser.get_text()[:40000]
            
        # Use AI to understand and extract the main content
        main_content = extract_main_content(raw_text)
        
        return {"text": main_content} 
    except Exception as e:
        return {"text": f"Error scraping URL: {str(e)}"}

@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    name = file.filename.lower() if file.filename else ""
    content = await file.read()
    text = ""
    try:
        if name.endswith(".pdf"):
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif name.endswith(".docx") or name.endswith(".doc"):
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        elif name.endswith(".pptx") or name.endswith(".ppt"):
            prs = pptx.Presentation(io.BytesIO(content))
            text_lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_lines.append(shape.text)
            text = "\n".join(text_lines)
        else:
            text = content.decode('utf-8')
    except Exception as e:
        text = f"Could not extract text from the file natively. Is it a valid document format?\nError: {str(e)}"
        
    return {"text": text[:40000]}

@router.get("/tts")
def proxy_tts(text: str, lang: str):
    import urllib.parse
    try:
        url = f"https://translate.googleapis.com/translate_tts?ie=UTF-8&q={urllib.parse.quote(text)}&tl={lang}&client=tw-ob"
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'
        }
        resp = requests.get(url, headers=headers, timeout=10)
        resp.raise_for_status()
        return Response(content=resp.content, media_type="audio/mpeg")
    except Exception as e:
        print(f"TTS Proxy Error: {e}")
        return Response(content=b"", status_code=500)